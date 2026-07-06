# -*- coding: utf-8 -*-
"""Generate MLOps_Soutenance.pptx — PowerPoint twin of index.html (25 slides).

Regenerate after swapping images in assets/:
    python make_pptx.py

python-pptx cannot script PowerPoint animations, so the deck is static;
speaker notes are included. Add entrance effects in PowerPoint if wanted
(select shapes -> Animations -> Fade).
"""

import math
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
A = lambda name: os.path.join(HERE, "assets", name)

# ── palette (light / white theme) ──────────────────────────────────────────
BG      = RGBColor(0xFF, 0xFF, 0xFF)
CARD    = RGBColor(0xF7, 0xF9, 0xFC)
LINE    = RGBColor(0xD7, 0xDE, 0xE8)
INK     = RGBColor(0x1F, 0x29, 0x37)
INK2    = RGBColor(0x47, 0x55, 0x69)
INK3    = RGBColor(0x64, 0x74, 0x8B)
CYAN    = RGBColor(0x0B, 0x53, 0x94)   # primary blue (readable on white)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
GOOD    = RGBColor(0x05, 0x96, 0x69)
BAD     = RGBColor(0xDC, 0x26, 0x26)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG  = RGBColor(0xF1, 0xF5, 0xF9)
# soft card tints against white
CYAN_T   = RGBColor(0xEA, 0xF4, 0xFB)
VIOLET_T = RGBColor(0xF3, 0xEF, 0xFC)
AMBER_T  = RGBColor(0xFD, 0xF3, 0xE3)
GOOD_T   = RGBColor(0xE8, 0xF7, 0xF0)
TEAL_T   = RGBColor(0xE6, 0xF5, 0xF5)
TEAL     = RGBColor(0x0F, 0x76, 0x6E)
CODE_TXT = RGBColor(0x33, 0x41, 0x55)
CODE_MUT = RGBColor(0x94, 0xA3, 0xB8)
CODE_GRN = RGBColor(0x04, 0x78, 0x57)

FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SLIDE_NO = 0


# ── helpers ────────────────────────────────────────────────────────────────
def new_slide(brand=True):
    global SLIDE_NO
    SLIDE_NO += 1
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    if brand:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.32), Inches(0.30), Inches(0.10), Inches(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = CYAN; dot.line.fill.background()
        dot.shadow.inherit = False
        tb(s, 0.48, 0.20, 3.4, 0.3, "MLOPS PLATFORM", size=10.5, color=INK3, bold=True, spacing=2.0)
    tb(s, 12.55, 7.08, 0.6, 0.3, str(SLIDE_NO), size=10, color=INK3, align=PP_ALIGN.RIGHT)
    return s


def tb(slide, x, y, w, h, text, size=16, color=INK2, bold=False, align=PP_ALIGN.LEFT,
       font=FONT, anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.name = font
        f.color.rgb = color
        if spacing:
            try:
                from pptx.oxml.ns import qn
                r._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
            except Exception:
                pass
    return box


def rich(slide, x, y, w, h, runs, size=16, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         font=FONT, line_spacing=None):
    """runs: list of paragraphs; each paragraph = list of (text, color, bold[, font]) tuples."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for item in para:
            text, color, bold = item[0], item[1], item[2]
            fnt = item[3] if len(item) > 3 else font
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.name = fnt
            r.font.color.rgb = color
    return box


def card(slide, x, y, w, h, fill=CARD, line=LINE, line_w=1.0, dash=None, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    if dash:
        sh.line.dash_style = dash
    sh.shadow.inherit = False
    sh.text_frame.margin_left = sh.text_frame.margin_right = Inches(0.08)
    sh.text_frame.margin_top = sh.text_frame.margin_bottom = Inches(0.04)
    return sh


def shape_text(sh, paras, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT
        r.font.color.rgb = color
    return sh


def pic(slide, name, x, y, w=None, h=None, frame=True):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    p = slide.shapes.add_picture(A(name), Inches(x), Inches(y), **kw)
    if frame:
        p.line.color.rgb = LINE
        p.line.width = Pt(1)
    p.shadow.inherit = False
    return p


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title(slide, plain, accent=""):
    rich(slide, 0.55, 0.62, 12.2, 0.7,
         [[(plain, INK, True), (accent, CYAN, True)]], size=27)


def checks(slide, items, x, y, w, col=2, size=13.5, row_h=0.52):
    per = math.ceil(len(items) / col)
    cw = w / col
    for i, item in enumerate(items):
        cx = x + (i // per) * cw
        cy = y + (i % per) * row_h
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.26), Inches(0.26))
        c.fill.solid(); c.fill.fore_color.rgb = GOOD_T
        c.line.color.rgb = GOOD; c.line.width = Pt(1); c.shadow.inherit = False
        shape_text(c, [("✓", 11, GOOD, True)])
        tb(slide, cx + 0.38, cy - 0.02, cw - 0.5, row_h, item, size=size, color=INK2)


def pill(slide, x, y, text, color=CYAN, tint=CYAN_T, w=1.35):
    sh = card(slide, x, y, w, 0.34, fill=tint, line=color, radius=0.5)
    shape_text(sh, [(text, 10.5, color, True)])
    return sh


def styled_table(slide, x, y, w, col_fracs, header, rows, row_h=0.62, size=12,
                 hl_last=False):
    cols = [w * f for f in col_fracs]
    # header
    cx = x
    for j, htext in enumerate(header):
        tb(slide, cx + 0.08, y, cols[j] - 0.16, 0.34, htext.upper(), size=10.5,
           color=INK3, bold=True)
        cx += cols[j]
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.36), Inches(w), Pt(1.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    # rows
    ry = y + 0.44
    for i, row in enumerate(rows):
        is_hl = hl_last and i == len(rows) - 1
        if is_hl:
            bgr = card(slide, x - 0.06, ry - 0.05, w + 0.12, row_h, fill=CYAN_T, line=CYAN, line_w=1.0, radius=0.12)
        cx = x
        for j, cell in enumerate(row):
            color = INK if (j == 0 or is_hl) else INK2
            bold = j == 0
            cw = cols[j] if j < len(cols) else w - (cx - x)
            tb(slide, cx + 0.08, ry, cw - 0.16, row_h, cell, size=size, color=color, bold=bold, line_spacing=0.95)
            cx += cw
        if not is_hl:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(ry + row_h - 0.07), Inches(w), Pt(0.75))
            sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(0xE5, 0xEA, 0xF1)
            sep.line.fill.background(); sep.shadow.inherit = False
        ry += row_h


def fnode(slide, x, y, w, h, lines, fill=CYAN_T, line=CYAN):
    sh = card(slide, x, y, w, h, fill=fill, line=line, radius=0.14)
    paras = [(lines[0], 12, INK, True)]
    for extra in lines[1:]:
        paras.append((extra, 9.5, INK3, False))
    shape_text(sh, paras)
    return sh


def farrow(slide, x, y, h=0.5):
    tb(slide, x, y + h / 2 - 0.16, 0.3, 0.32, "→", size=15, color=CYAN, bold=True, align=PP_ALIGN.CENTER)


def code_box(slide, x, y, w, h, paras, size=11):
    sh = card(slide, x, y, w, h, fill=CODEBG, line=LINE, radius=0.06)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        for text, color in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.name = MONO
            r.font.color.rgb = color
    return sh


# ═══ 1 — TITLE ═════════════════════════════════════════════════════════════
s = new_slide(brand=False)
pic(s, "esprit.png", 0.55, 0.45, h=0.62, frame=False)
pic(s, "insomea.png", 11.05, 0.45, h=0.58, frame=False)
p = pill(s, 5.42, 1.62, "END-OF-STUDIES PROJECT · 2025 / 2026", w=2.5)
tb(s, 0.8, 2.15, 11.73, 1.6,
   "Design & Implementation of an\nEnd-to-End MLOps Platform",
   size=37, color=INK, bold=True, align=PP_ALIGN.CENTER)
tb(s, 0.8, 3.62, 11.73, 0.5, "Automated Model Training, Tracking & Deployment on Kubernetes",
   size=17, color=INK2, align=PP_ALIGN.CENTER)
bar = slide_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.92), Inches(4.35), Inches(1.5), Pt(3.2))
bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background(); bar.shadow.inherit = False
who = [("PRESENTED BY", "Firas Mahjoubi"), ("ACADEMIC SUPERVISOR", "Mr. Ben Mardes Achref"),
       ("COMPANY SUPERVISOR", "Mr. Amine Gonji · INSOMEA")]
wx = 1.5
for label, name in who:
    tb(s, wx, 5.05, 3.5, 0.3, label, size=10.5, color=INK3, bold=True)
    tb(s, wx, 5.35, 3.5, 0.4, name, size=14.5, color=INK2)
    wx += 3.6
notes(s, "Good morning. Today I present my end-of-studies project, carried out at INSOMEA: a complete MLOps platform that takes a model from a notebook to a live prediction service.")

# ═══ 2 — AGENDA ════════════════════════════════════════════════════════════
s = new_slide()
title(s, "Agenda")
items = ["Host company", "Project context", "Problem & objectives", "State of the art",
         "Proposed solution", "Business & system objectives", "Methodology", "Technology choices",
         "Architecture & modeling", "The platform, sprint by sprint", "Case study & live demo",
         "Conclusion & perspectives"]
for i, it in enumerate(items):
    cx = 0.9 + (i // 6) * 6.1
    cy = 1.75 + (i % 6) * 0.78
    tb(s, cx, cy, 0.5, 0.4, f"{i+1:02d}", size=15, color=CYAN, bold=True, font=MONO)
    tb(s, cx + 0.62, cy + 0.01, 5.2, 0.4, it, size=16, color=INK2)
notes(s, "Quick roadmap: company, context and problem; then the solution, its architecture, the four sprints, a real case study, and the live demo.")

# ═══ 3 — HOST COMPANY ══════════════════════════════════════════════════════
s = new_slide()
title(s, "Host Company — ", "INSOMEA")
tb(s, 0.55, 1.42, 12.2, 0.4,
   "Microsoft Cloud Gold Partner · Middle East & Africa · 6 offices · 500+ customers · Microsoft Partner of the Year ×5",
   size=13, color=INK3)
pillars = [("Modern Workplace", "Collaboration & productivity with Microsoft 365 and Teams — adoption & change management.", CYAN),
           ("Security", "Identity & access management, information protection, threat protection.", VIOLET),
           ("Cloud & Infrastructure", "Microsoft Azure, datacenter, DevOps, application development & integration — where this project lives.", AMBER)]
for i, (h, ptext, color) in enumerate(pillars):
    cardsh = card(s, 0.55 + i * 4.18, 2.15, 3.9, 2.3)
    tb(s, 0.8 + i * 4.18, 2.4, 3.4, 0.4, h, size=16, color=color, bold=True)
    tb(s, 0.8 + i * 4.18, 2.9, 3.4, 1.4, ptext, size=12.5, color=INK2)
pic(s, "insomea.png", 5.2, 5.1, h=0.55, frame=False)
tb(s, 0.55, 6.0, 12.2, 0.4, "This PFE was conducted in the Cloud & Infrastructure team over 6 months.",
   size=13.5, color=INK2, align=PP_ALIGN.CENTER)
notes(s, "INSOMEA is a Microsoft Gold Partner operating across MEA. Three pillars; my project belongs to Cloud & Infrastructure — building an MLOps stack on Azure.")

# ═══ 4 — CONTEXT ═══════════════════════════════════════════════════════════
s = new_slide()
title(s, "Project Context")
ctx = [("Machine learning is stuck in notebooks",
        "Models are trained on laptops, results tracked in spreadsheets, deployment is manual and fragile — the classic “level 0” workflow.", CYAN),
       ("Industrialize the full ML lifecycle",
        "One platform: upload code → train on a Kubernetes cluster → track every experiment → register, deploy & consume models.", VIOLET),
       ("Deployed on real cloud infrastructure",
        "Runs in production on Azure Kubernetes Service, continuously delivered by CI/CD, publicly reachable behind Cloudflare.", AMBER)]
for i, (h, ptext, color) in enumerate(ctx):
    y = 1.7 + i * 1.72
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y), Pt(4), Inches(1.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background(); bar.shadow.inherit = False
    card(s, 0.65, y, 12.0, 1.45)
    tb(s, 0.95, y + 0.18, 11.4, 0.4, h, size=17, color=INK, bold=True)
    tb(s, 0.95, y + 0.62, 11.4, 0.7, ptext, size=13.5, color=INK2)
notes(s, "The gap between a model that works on a laptop and one that serves production traffic is the hardest problem in applied ML — that gap is exactly what the platform closes.")

# ═══ 5 — PROBLEM & OBJECTIVES ══════════════════════════════════════════════
s = new_slide()
title(s, "Problem & Objectives")
pb = card(s, 0.9, 1.6, 11.5, 1.05, fill=CYAN_T, line=CYAN, radius=0.12)
shape_text(pb, [("How can a data scientist go from a notebook to a live, consumable prediction service —", 16.5, INK, True),
                ("without writing a single line of infrastructure code?", 16.5, CYAN, True)])
tb(s, 0.9, 3.05, 6, 0.4, "Objectives of the PFE:", size=14, color=INK, bold=True)
checks(s, ["Upload & run any code on the cluster", "Track every experiment automatically",
           "Registry: version & promote models", "One-click deployment as live services",
           "Public API for external applications", "Monitoring, diagnostics & administration"],
       0.9, 3.6, 11.6, col=2, size=14.5, row_h=0.78)
notes(s, "One question drives everything. Six objectives: run, track, register, deploy, expose, operate. All six were delivered.")

# ═══ 6 — STATE OF THE ART ══════════════════════════════════════════════════
s = new_slide()
title(s, "State of the Art")
styled_table(
    s, 0.6, 1.6, 12.1, [0.24, 0.33, 0.43],
    ["Solution", "Strengths", "Limitations"],
    [
        ["SageMaker / Vertex AI /\nAzure ML", "Full managed lifecycle", "Vendor lock-in, recurring cost, hidden mechanics"],
        ["Databricks", "Strong data + training environment", "Subscription; data-centric, not serving-centric"],
        ["AIchor (InstaDeep)", "Cluster training via GitOps", "Stops at training — no registry, no serving, no public API; needs Git + Docker skills"],
        ["Kubeflow (open source)", "Powerful primitives", "No product: no users, no projects, heavy integration"],
        ["This project", "Full lifecycle — training → tracking → registry → serving → public API — 100% open-source, self-hosted, browser-only for the user", ""],
    ],
    row_h=0.86, size=12.5, hl_last=True,
)
notes(s, "Managed clouds are complete but locked and opaque. AIchor is the closest comparable but stops at training. Our platform covers the full loop, open source, and asks nothing from the user but a browser.")

# ═══ 7 — LIFECYCLE RING ════════════════════════════════════════════════════
s = new_slide()
title(s, "Proposed Solution — ", "the complete ML lifecycle, one platform")
cx, cy, rx, ry = 6.67, 4.25, 3.55, 2.15
orbit = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - rx + 0.6), Inches(cy - ry + 0.28),
                           Inches((rx - 0.6) * 2), Inches((ry - 0.28) * 2))
orbit.fill.background(); orbit.line.color.rgb = RGBColor(0xB9, 0xCB, 0xDD)
orbit.line.width = Pt(1.5); orbit.line.dash_style = MSO_LINE.DASH; orbit.shadow.inherit = False
tb(s, cx - 1.5, cy - 0.35, 3.0, 0.8, "The end-to-end\nMLOps lifecycle", size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
ring_nodes = [
    ("Upload", "code + data", CYAN_T, CYAN), ("Train", "isolated pod", CYAN_T, CYAN),
    ("Track", "autolog", CYAN_T, CYAN), ("Register", "versions", VIOLET_T, VIOLET),
    ("Deploy", "KServe", VIOLET_T, VIOLET), ("Serve", "predictions", VIOLET_T, VIOLET),
    ("Consume", "public API", AMBER_T, AMBER), ("Monitor", "KPIs · telemetry", AMBER_T, AMBER),
]
NW, NH = 1.6, 0.72
for i, (label, sub, fill, line) in enumerate(ring_nodes):
    ang = math.radians(i * 45)  # 0 = top, clockwise
    nx = cx + rx * math.sin(ang) - NW / 2
    ny = cy - ry * math.cos(ang) - NH / 2
    n = card(s, nx, ny, NW, NH, fill=fill, line=line, radius=0.18)
    shape_text(n, [(label, 12.5, INK, True), (sub, 9.5, INK3, False)])
notes(s, "Walk the ring stage by stage — this IS the demo path. Upload, train, track, register, deploy, serve, consume from outside, monitor. Every arrow is a platform feature, not a manual step.")

# ═══ 8 — BO & DSO ══════════════════════════════════════════════════════════
s = new_slide()
title(s, "Business & System Objectives")
styled_table(
    s, 0.6, 1.55, 12.1, [0.42, 0.58],
    ["Business objective", "Delivered system capability"],
    [
        ["Train without owning infrastructure", "Any .py / .ipynb / .zip runs in an isolated, resource-limited Kubernetes pod"],
        ["Never lose an experiment", "MLflow autolog captures params, metrics & artifacts of every run — zero code changes"],
        ["Choose the best model objectively", "Leaderboard groups runs by model family, ranked on evaluation metrics only"],
        ["Ship a model like software", "Registry versions + stages; one-click deployment as a KServe InferenceService"],
        ["Let external apps consume models", "Public endpoint secured by hashed API keys — no platform account needed"],
        ["Operate & trust the system", "Model-evolution KPIs, serving telemetry, failure diagnostics, admin area"],
    ],
    row_h=0.8, size=12.5,
)
notes(s, "Each business need maps to a concrete shipped capability — this table is the contract the jury can check against the demo.")

# ═══ 9 — METHODOLOGY ═══════════════════════════════════════════════════════
s = new_slide()
title(s, "Methodology — ", "Scrum, 4 sprints")
tb(s, 0.55, 1.42, 12.2, 0.4,
   "6 months (12 Jan → 11 Jul 2026) · product backlog of 15 user stories · a demonstrable increment per sprint",
   size=13, color=INK3)
sprints = [
    ("SPRINT 1", "Foundation", "JWT auth, projects, one MLflow experiment per project", "02 Feb → 01 Mar", CYAN),
    ("SPRINT 2", "Training & Tracking", "Upload, pipeline on K8s, notebook conversion, live logs, autolog", "02 Mar → 05 Apr", VIOLET),
    ("SPRINT 3", "Registry & Deployment", "Model versions & stages, KServe serving, prediction tester", "06 Apr → 03 May", AMBER),
    ("SPRINT 4", "Observability & Public Serving", "Dashboards, public API keys, diagnostics, administration", "04 May → 14 Jun", GOOD),
]
for i, (n, t, d, dates, color) in enumerate(sprints):
    x = 0.55 + i * 3.13
    topbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.15), Inches(2.95), Pt(4))
    topbar.fill.solid(); topbar.fill.fore_color.rgb = color; topbar.line.fill.background(); topbar.shadow.inherit = False
    card(s, x, 2.2, 2.95, 2.9)
    tb(s, x + 0.2, 2.42, 2.55, 0.3, n, size=11, color=INK3, bold=True, spacing=1.5)
    tb(s, x + 0.2, 2.76, 2.55, 0.75, t, size=15.5, color=INK, bold=True)
    tb(s, x + 0.2, 3.55, 2.55, 1.1, d, size=11.5, color=INK2)
    tb(s, x + 0.2, 4.72, 2.55, 0.3, dates, size=10.5, color=INK3, font=MONO)
tb(s, 0.55, 5.6, 12.2, 0.4,
   "Product Owner: Mr. A. Gonji (INSOMEA) · Scrum Master: Mr. B. M. Achref · Development team: the student",
   size=12.5, color=INK3, align=PP_ALIGN.CENTER)
notes(s, "Scrum with four vertical slices — each sprint ended with something running and reviewed. This order is deliberate: identity → training → serving → operations.")

# ═══ 10 — USE CASE ═════════════════════════════════════════════════════════
s = new_slide()
title(s, "Global Use Case Diagram")
panel = card(s, 2.35, 1.5, 8.63, 5.6, fill=WHITE, line=LINE, radius=0.04)
pic(s, "usecase_global.png", 2.55, 1.65, h=5.3, frame=False)
notes(s, "Three human actors: the data scientist (whole lifecycle), the external API client (one authenticated HTTP call), and the administrator (accounts, oversight).")

# ═══ 11 — TECH STACK ═══════════════════════════════════════════════════════
s = new_slide()
title(s, "Technology Choices")
cols = [
    ("MLOPS CORE", CYAN, [("mlflow.png", "MLflow — tracking & registry"),
                          ("kserve.png", "KServe — model serving"),
                          ("kubeflow.png", "Kubeflow Pipelines + Argo"),
                          ("minio.png", "MinIO — S3 object storage")]),
    ("APPLICATION", VIOLET, [("angular.png", "Angular 19 — frontend SPA"),
                             ("fastapi.png", "FastAPI — async backend"),
                             ("postgresql.png", "PostgreSQL — platform + MLflow")]),
    ("INFRASTRUCTURE", AMBER, [("kubernetes.png", "Kubernetes — AKS / KinD"),
                               ("github_actions.png", "GitHub Actions — CI/CD"),
                               ("acr.png", "Azure Container Registry"),
                               ("cloudflare.png", "Cloudflare Tunnel — secure edge")]),
]
for i, (h, color, items) in enumerate(cols):
    x = 0.55 + i * 4.18
    tb(s, x, 1.55, 3.9, 0.35, h, size=13.5, color=color, bold=True, spacing=1.5)
    for j, (logo, label) in enumerate(items):
        y = 2.05 + j * 0.92
        card(s, x, y, 3.9, 0.78)
        pic(s, logo, x + 0.18, y + 0.14, h=0.5, frame=False)
        tb(s, x + 1.0, y + 0.24, 2.85, 0.5, label, size=12, color=INK, bold=True)
notes(s, "Three families: the MLOps engines, the application around them, and the infrastructure. Everything open source; Azure is the ground it runs on.")

# ═══ 12 — LOGICAL ARCHITECTURE ═════════════════════════════════════════════
s = new_slide()
title(s, "Logical Architecture")
layers = [
    ("Presentation — Angular SPA", "dashboard · projects · experiments · deployments · admin"),
    ("API — FastAPI routers", "validation · authentication · ownership checks"),
    ("Service — business logic", "pipeline submission · MLflow · KServe · storage orchestration"),
    ("Data — SQLAlchemy / PostgreSQL", "users · projects · runs · models · deployments · API keys · telemetry"),
]
for i, (h, sub) in enumerate(layers):
    y = 1.6 + i * 1.32
    lay = card(s, 0.55, y, 7.4, 1.0, fill=CYAN_T, line=CYAN, radius=0.12)
    tb(s, 0.85, y + 0.14, 6.9, 0.4, h, size=15, color=INK, bold=True)
    tb(s, 0.85, y + 0.55, 6.9, 0.35, sub, size=11, color=INK2)
    if i < 3:
        tb(s, 4.05, y + 1.0, 0.5, 0.3, "▼", size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
sysbox = card(s, 8.45, 1.9, 4.3, 4.2, fill=BG, line=AMBER, dash=MSO_LINE.DASH, radius=0.1)
tb(s, 8.7, 1.72, 3.0, 0.3, "CLUSTER SYSTEMS", size=10.5, color=AMBER, bold=True, spacing=1.5)
sysrows = [("mlflow.png", "MLflow — experiments & registry"), ("minio.png", "MinIO — code, datasets, artifacts"),
           ("kserve.png", "KServe — InferenceServices"), ("kubeflow.png", "KFP / Argo — training pods")]
for j, (logo, label) in enumerate(sysrows):
    y = 2.25 + j * 0.92
    card(s, 8.65, y, 3.9, 0.75, fill=AMBER_T, line=AMBER, radius=0.14)
    pic(s, logo, 8.82, y + 0.14, h=0.46, frame=False)
    tb(s, 9.6, y + 0.2, 2.9, 0.5, label, size=11, color=INK, bold=True)
notes(s, "Classic layered backend — deliberately NOT microservices: the heavy concerns already live in dedicated cluster systems shown on the right; the service layer orchestrates them.")

# ═══ 13 — DEPLOYMENT ARCHITECTURE ══════════════════════════════════════════
s = new_slide()
title(s, "Deployment Architecture")
card(s, 2.0, 1.42, 9.33, 5.35, fill=WHITE, line=LINE, radius=0.04)
pic(s, "architecture.png", 2.2, 1.57, h=5.05, frame=False)
tb(s, 0.55, 6.9, 12.2, 0.4,
   "browser → Cloudflare edge (TLS) → tunnel → frontend → backend → { PostgreSQL · MinIO · MLflow · KFP/Argo · KServe } — all on AKS",
   size=11.5, color=INK3, align=PP_ALIGN.CENTER, font=MONO)
notes(s, "Everything is containerized on AKS. No public IP is exposed: an outbound-only Cloudflare Tunnel fronts the cluster, TLS terminates at the edge, DDoS protection included.")

# ═══ 14 — SPRINT 1 ═════════════════════════════════════════════════════════
s = new_slide()
pill(s, 0.55, 0.66, "SPRINT 1", w=1.1)
rich(s, 1.85, 0.62, 11, 0.6, [[("Foundation — ", INK, True), ("identity & projects", CYAN, True)]], size=27)
pic(s, "screen_login.png", 0.55, 1.7, w=6.05)
tb(s, 0.55, 5.2, 6.05, 0.6, "Stateless JWT auth — access + refresh tokens, bcrypt hashing", size=11.5, color=INK3)
pic(s, "screen_projects.png", 6.85, 1.7, w=6.05)
tb(s, 6.85, 5.2, 6.05, 0.6, "Projects — each bound to its own MLflow experiment; strict ownership (404 for foreign resources)", size=11.5, color=INK3)
notes(s, "Sprint 1 delivers identity and the project as organizing unit. Creating a project provisions its MLflow experiment — that binding powers everything later.")

# ═══ 15 — SPRINT 2 PIPELINE ════════════════════════════════════════════════
s = new_slide()
pill(s, 0.55, 0.66, "SPRINT 2", w=1.1)
rich(s, 1.85, 0.62, 11, 0.6, [[("Training & Tracking — ", INK, True), ("the engine", CYAN, True)]], size=27)
fy = 1.75
fnode(s, 0.55, fy, 1.25, 0.85, ["Trigger", "Run code"])
farrow(s, 1.83, fy, 0.85)
podbox = card(s, 2.2, fy - 0.28, 7.35, 1.42, fill=BG, line=CYAN, dash=MSO_LINE.DASH, radius=0.1)
tb(s, 2.45, fy - 0.47, 4.0, 0.3, "ISOLATED POD — KUBERNETES", size=9.5, color=CYAN, bold=True, spacing=1.2)
steps = [("Download", "code + data"), ("Install", "dependencies"), ("Autolog", "runner wraps"), ("Execute", "auto-retry deps")]
sx = 2.45
for i, (h, sub) in enumerate(steps):
    fnode(s, sx, fy, 1.45, 0.85, [h, sub], fill=TEAL_T, line=TEAL)
    sx += 1.5
    if i < 3:
        farrow(s, sx - 0.06, fy, 0.85)
        sx += 0.28
farrow(s, 9.62, fy, 0.85)
fnode(s, 10.0, fy, 1.7, 0.85, ["MLflow", "metrics · params · model"], fill=AMBER_T, line=AMBER)
pic(s, "screen_logs.png", 2.8, 3.35, w=7.7)
tb(s, 2.8, 6.85, 7.7, 0.4, "Live pod logs streamed to the browser during the run", size=11.5, color=INK3, align=PP_ALIGN.CENTER)
notes(s, "A run is a small pipeline in an isolated pod: download, install, wrap with the autolog runner, execute with automatic dependency retries. The user watches live logs; MLflow captures everything with zero code changes.")

# ═══ 16 — NOTEBOOK CONVERSION ══════════════════════════════════════════════
s = new_slide()
title(s, "Signature Feature — ", "notebook → script conversion")
tb(s, 0.55, 1.4, 12.2, 0.6,
   "Real-world notebooks are full of Jupyter-only constructs. The platform converts them to clean scripts at upload time — using IPython's own transformer, not fragile regexes.",
   size=12.5, color=INK3)
code_box(s, 0.55, 2.25, 5.7, 2.3, [
    [("# messy Colab notebook", INK3)],
    [("!pip install", BAD), (" xgboost", INK2)],
    [("%matplotlib", BAD), (" inline", INK2)],
    [("%%time", BAD)],
    [("model.fit(X_train, y_train)", INK2)],
    [("files = ", INK2), ("!ls", BAD)],
], size=12)
tb(s, 6.35, 3.1, 0.6, 0.5, "→", size=26, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
code_box(s, 7.05, 2.25, 5.7, 2.3, [
    [("# clean runnable script", INK3)],
    [("subprocess.check_call([sys.executable,", GOOD)],
    [("  '-m','pip','install','xgboost'])", GOOD)],
    [("# magic removed — code kept:", INK3)],
    [("model.fit(X_train, y_train)", INK2)],
    [("files = ", INK2), ("subprocess.run('ls', ...)", GOOD)],
], size=12)
checks(s, ["String-safe — magics inside strings untouched", "!pip install preserved as real installs",
           "%%time body kept — training never lost", "Output guaranteed to compile, or safe fallback"],
       0.9, 5.05, 11.6, col=2, size=13.5, row_h=0.62)
notes(s, "This is the feature that makes 'upload any notebook' true. Cells that already parse are untouched; magic cells go through IPython's canonical transformer, then get rewritten into portable Python. Warnings surface instantly in the UI.")

# ═══ 17 — EXPERIMENTS INTELLIGENCE ═════════════════════════════════════════
s = new_slide()
title(s, "Experiments Intelligence")
pic(s, "screen_experiments.png", 1.9, 1.45, w=9.5)
checks(s, ["Runs grouped by model family (autolog tags)", "Ranked on evaluation metrics only — never training scores",
           "Champion banner + star ratings", "Register any run as a model version, in one click"],
       0.9, 6.15, 11.8, col=2, size=12.5, row_h=0.52)
notes(s, "One notebook execution can produce many runs — the leaderboard buckets them per algorithm and crowns a champion using eval metrics only, so overfit training scores can't cheat. From here, one click registers the best run.")

# ═══ 18 — SPRINT 3 DEPLOY FLOW ═════════════════════════════════════════════
s = new_slide()
pill(s, 0.55, 0.66, "SPRINT 3", w=1.1)
rich(s, 1.85, 0.62, 11, 0.6, [[("Registry & Deployment", INK, True)]], size=27)
fy = 1.7
flow = [
    (["Register", "model version"], CYAN_T, CYAN), (["InferenceService", "one K8s resource"], CYAN_T, CYAN),
    (["KServe operator", "reconciles"], TEAL_T, TEAL),
    (["Deployment", "Service · Ingress"], TEAL_T, TEAL),
    (["Model server", "loads from MinIO"], AMBER_T, AMBER), (["READY", "predict!"], GOOD_T, GOOD),
]
fx = 0.55
for i, (lines, fill, line) in enumerate(flow):
    w = 1.85
    fnode(s, fx, fy, w, 0.9, lines, fill=fill, line=line)
    fx += w
    if i < len(flow) - 1:
        farrow(s, fx - 0.03, fy, 0.9)
        fx += 0.28
pic(s, "screen_registry.png", 0.55, 3.15, w=6.05)
tb(s, 0.55, 6.6, 6.05, 0.6, "Registry — versions & lifecycle stages (Staging / Production / Archived)", size=11, color=INK3)
pic(s, "screen_deploy.png", 6.85, 3.15, w=6.05)
tb(s, 6.85, 6.6, 6.05, 0.6, "Deployment panel — status progress & endpoint", size=11, color=INK3)
notes(s, "Deployment is declarative: the platform creates ONE Kubernetes resource, the KServe operator does the rest — deployment, service, ingress, model download. The UI tracks it to READY.")

# ═══ 19 — PREDICTION & PUBLIC API ══════════════════════════════════════════
s = new_slide()
title(s, "Prediction & ", "Public API")
pic(s, "screen_predict.png", 0.55, 1.6, w=6.6)
tb(s, 0.55, 5.5, 6.6, 0.4, "In-platform prediction tester", size=11.5, color=INK3)
code_box(s, 7.5, 1.6, 5.3, 2.5, [
    [("curl -X POST ", INK2), ("https://mlops.firasmahjoubi.app", CYAN)],
    [("  /api/public/predict/<deployment>", CYAN)],
    [("  -H ", INK2), ('"Authorization: Bearer mlops_****"', GOOD)],
    [("  -H \"Content-Type: application/json\"", INK2)],
    [("  -d ", INK2), ("'{\"instances\": [[17.9, 10.3, ...]]}'", GOOD)],
    [("", INK2)],
    [("→ {\"predictions\": [0]}   # 45 ms", INK3)],
], size=11)
checks(s, ["Per-deployment API keys — SHA-256 hashed, shown once",
           "Any website / app can consume the model — no account",
           "Ready-to-paste cURL / JS / Python snippets generated"],
       7.5, 4.4, 5.3, col=1, size=12, row_h=0.62)
notes(s, "Deployed models are consumable from outside through a key-authenticated public endpoint. Keys are stored hashed, revocable, and the UI generates the integration snippets.")

# ═══ 20 — SPRINT 4 ═════════════════════════════════════════════════════════
s = new_slide()
pill(s, 0.55, 0.66, "SPRINT 4", w=1.1)
rich(s, 1.85, 0.62, 11, 0.6, [[("Observability & Administration", INK, True)]], size=27)
pic(s, "screen_dashboard.png", 0.55, 1.6, w=6.05)
tb(s, 0.55, 4.32, 6.05, 0.35, "Live cluster & per-pod metrics", size=11, color=INK3)
pic(s, "screen_admin.png", 6.85, 1.6, w=6.05)
tb(s, 6.85, 4.32, 6.05, 0.35, "Administration — users, roles, platform oversight", size=11, color=INK3)
pic(s, "screen_diagnostics.png", 0.55, 4.85, w=6.05)
tb(s, 0.55, 6.98, 6.05, 0.35, "Failure diagnostics — “why did my run fail?” at a glance", size=11, color=INK3)
mon = card(s, 6.85, 4.85, 6.05, 1.95, fill=CARD, line=LINE)
tb(s, 7.1, 5.05, 5.6, 0.35, "+ Per-project Monitoring", size=14.5, color=CYAN, bold=True)
tb(s, 7.1, 5.45, 5.6, 1.3,
   "Model evolution over time per family · health verdict (improving / degrading) · version-to-version deltas · serving telemetry: request volume, latency p95, error rate — from real prediction logs.",
   size=11.5, color=INK2)
notes(s, "Sprint 4 makes the platform operable: metrics dashboards, admin area, self-explaining failures, and per-project monitoring — is my model getting better or worse, and how is it behaving in production?")

# ═══ 21 — CASE STUDY ═══════════════════════════════════════════════════════
s = new_slide()
title(s, "Case Study — ", "telecom churn, end to end")
tb(s, 0.55, 1.42, 12.2, 0.4,
   "A real churn model (XGBoost + SMOTE, built from 2 Excel sources) uploaded as a notebook zip — the platform did the rest.",
   size=13, color=INK3)
tiles = [("6", "TRACKED RUNS\nFROM 1 UPLOAD", CYAN), ("0.91", "ROC-AUC\nTEST SET", GOOD),
         ("0.85", "F1 SCORE\nCHURN CLASS", GOOD), ("96%", "CHURN CAPTURED IN\nHIGH-RISK SEGMENT", AMBER)]
for i, (v, l, color) in enumerate(tiles):
    x = 1.0 + i * 2.95
    card(s, x, 2.1, 2.6, 1.75)
    tb(s, x, 2.35, 2.6, 0.7, v, size=34, color=color, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x, 3.15, 2.6, 0.6, l, size=9.5, color=INK3, bold=True, align=PP_ALIGN.CENTER)
checks(s, ["Notebook auto-converted (pip installs preserved)", "3 model families ranked on the leaderboard",
           "Champion registered → deployed → public API", "Monitored: evolution, latency, error rate"],
       1.0, 4.45, 11.4, col=2, size=14, row_h=0.75)
notes(s, "Full-circle proof: a colleague's real DW churn model, adapted to its two Excel sources, uploaded as-is. Six runs tracked, XGBoost champion at 0.91 AUC, deployed and consumed through the public API — every platform feature exercised by one real use case.")

# ═══ 22 — DEMO ═════════════════════════════════════════════════════════════
s = new_slide()
tb(s, 0.8, 2.5, 11.73, 1.3, "LIVE DEMO", size=64, color=CYAN, bold=True, align=PP_ALIGN.CENTER, spacing=3.0)
tb(s, 0.8, 4.1, 11.73, 0.6,
   "upload a notebook  →  watch it train  →  leaderboard  →  register the champion  →  deploy  →  predict from the public API",
   size=14, color=INK2, align=PP_ALIGN.CENTER)
notes(s, "Switch to the browser: mlops.firasmahjoubi.app. Follow the ring from slide 7 live. Fallback if network fails: screenshots + the recorded run.")

# ═══ 23 — CONCLUSION ═══════════════════════════════════════════════════════
s = new_slide()
title(s, "Conclusion")
conc = [("End-to-end lifecycle, automated",
         "Notebook → tracked runs → registry → live service → public API — from a single web interface, no infra knowledge required.", CYAN),
        ("MLOps maturity: level 0 → 1",
         "Repeatable pipelines, versioned models, one-click deployment; the platform itself delivered at level 2 through full CI/CD.", VIOLET),
        ("Real, production-grade infrastructure",
         "Runs on AKS behind Cloudflare, fully open-source stack, validated by a real churn use case.", GOOD)]
for i, (h, ptext, color) in enumerate(conc):
    x = 0.55 + i * 4.18
    topbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(3.9), Pt(4))
    topbar.fill.solid(); topbar.fill.fore_color.rgb = color; topbar.line.fill.background(); topbar.shadow.inherit = False
    card(s, x, 2.05, 3.9, 2.7)
    tb(s, x + 0.25, 2.3, 3.4, 0.75, h, size=15.5, color=INK, bold=True)
    tb(s, x + 0.25, 3.1, 3.4, 1.5, ptext, size=12, color=INK2)
notes(s, "Three takeaways: the loop is closed, the maturity jump is measurable, and it's not a prototype — it runs on real cloud infrastructure today.")

# ═══ 24 — PERSPECTIVES ═════════════════════════════════════════════════════
s = new_slide()
title(s, "Perspectives")
checks(s, ["Data-drift detection — feature stats at serve time", "Alerting — thresholds on error rate & degradation",
           "A/B testing & canary between model versions", "GPU scheduling for deep-learning workloads",
           "Multi-tenancy — teams & role-based access", "Curated training image — pre-pinned ML packages"],
       1.0, 2.0, 11.5, col=2, size=15.5, row_h=1.05)
notes(s, "The monitoring foundation added in sprint 4 makes drift detection and alerting natural next steps; GPU and multi-tenancy would take it from single-team to organization-wide.")

# ═══ 25 — THANK YOU ════════════════════════════════════════════════════════
s = new_slide(brand=False)
pic(s, "esprit.png", 0.55, 0.45, h=0.62, frame=False)
pic(s, "insomea.png", 11.05, 0.45, h=0.58, frame=False)
tb(s, 0.8, 2.7, 11.73, 1.2, "Thank you!", size=54, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.92), Inches(4.15), Inches(1.5), Pt(3.2))
bar.fill.solid(); bar.fill.fore_color.rgb = VIOLET; bar.line.fill.background(); bar.shadow.inherit = False
tb(s, 0.8, 4.55, 11.73, 0.5, "Questions are welcome · Firas Mahjoubi · mlops.firasmahjoubi.app",
   size=14, color=INK3, align=PP_ALIGN.CENTER)
notes(s, "Thank the jury; invite questions. Have the platform tab still open for follow-up questions.")

# ── save ───────────────────────────────────────────────────────────────────
OUT = os.path.join(HERE, sys.argv[1] if len(sys.argv) > 1 else "MLOps_Soutenance.pptx")
prs.save(OUT)
print(f"wrote {OUT} ({os.path.getsize(OUT) / 1e6:.1f} MB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
